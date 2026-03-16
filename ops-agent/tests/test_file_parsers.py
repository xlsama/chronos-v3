"""Tests for file parsers."""

import pytest

from src.lib.file_parsers import (
    IMAGE_EXTENSIONS,
    SUPPORTED_EXTENSIONS,
    ParsedSegment,
    is_image,
    parse_csv,
    parse_file,
    parse_file_segments,
    parse_text,
)


class TestParseText:
    def test_parses_utf8(self):
        content = "Hello, world!".encode("utf-8")
        assert parse_text(content) == "Hello, world!"

    def test_handles_non_utf8(self):
        content = b"\xff\xfe"
        result = parse_text(content)
        assert isinstance(result, str)


class TestParseCsv:
    def test_parses_simple_csv(self):
        content = "name,age\nAlice,30\nBob,25".encode("utf-8")
        result = parse_csv(content)
        assert "name" in result
        assert "Alice" in result
        assert "Bob" in result
        assert "|" in result  # markdown table format

    def test_empty_csv(self):
        content = b""
        assert parse_csv(content) == ""


class TestParseFile:
    def test_text_file(self):
        content = b"Hello text"
        assert parse_file(content, "readme.txt") == "Hello text"

    def test_markdown_file(self):
        content = b"# Title"
        assert parse_file(content, "readme.md") == "# Title"

    def test_csv_file(self):
        content = b"a,b\n1,2"
        result = parse_file(content, "data.csv")
        assert "a" in result
        assert "1" in result

    def test_unsupported_format(self):
        with pytest.raises(ValueError, match="Unsupported"):
            parse_file(b"data", "file.xyz")


class TestSupportedExtensions:
    def test_has_common_formats(self):
        assert ".pdf" in SUPPORTED_EXTENSIONS
        assert ".docx" in SUPPORTED_EXTENSIONS
        assert ".xlsx" in SUPPORTED_EXTENSIONS
        assert ".csv" in SUPPORTED_EXTENSIONS
        assert ".md" in SUPPORTED_EXTENSIONS
        assert ".txt" in SUPPORTED_EXTENSIONS

    def test_has_new_formats(self):
        assert ".pptx" in SUPPORTED_EXTENSIONS
        assert ".html" in SUPPORTED_EXTENSIONS
        assert ".htm" in SUPPORTED_EXTENSIONS
        assert ".json" in SUPPORTED_EXTENSIONS
        assert ".yaml" in SUPPORTED_EXTENSIONS
        assert ".yml" in SUPPORTED_EXTENSIONS
        assert ".log" in SUPPORTED_EXTENSIONS

    def test_includes_image_extensions(self):
        for ext in IMAGE_EXTENSIONS:
            assert ext in SUPPORTED_EXTENSIONS


class TestIsImage:
    def test_png(self):
        assert is_image("photo.png") is True

    def test_jpg(self):
        assert is_image("photo.jpg") is True

    def test_jpeg(self):
        assert is_image("photo.jpeg") is True

    def test_webp(self):
        assert is_image("photo.webp") is True

    def test_case_insensitive(self):
        assert is_image("photo.PNG") is True

    def test_non_image(self):
        assert is_image("doc.pdf") is False
        assert is_image("text.txt") is False


class TestParsePptx:
    def test_parse_pptx(self):
        """parse_pptx returns text from slides."""
        from pptx import Presentation
        import io

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank layout
        txBox = slide.shapes.add_textbox(0, 0, 100, 100)
        txBox.text_frame.paragraphs[0].text = "Hello from slide 1"

        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        txBox2 = slide2.shapes.add_textbox(0, 0, 100, 100)
        txBox2.text_frame.paragraphs[0].text = "Hello from slide 2"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        result = parse_file(buf.read(), "test.pptx")
        assert "Hello from slide 1" in result
        assert "Hello from slide 2" in result


class TestParseHtml:
    def test_parse_html(self):
        html = b"<html><body><h1>Title</h1><p>Content</p><script>evil()</script></body></html>"
        result = parse_file(html, "page.html")
        assert "Title" in result
        assert "Content" in result
        assert "evil" not in result


class TestParseJson:
    def test_parse_json(self):
        data = b'{"key": "value", "num": 42}'
        result = parse_file(data, "config.json")
        assert "key" in result
        assert "value" in result
        assert "42" in result


class TestParseFileSegments:
    def test_text_returns_single_segment(self):
        segments = parse_file_segments(b"Hello world", "test.txt")
        assert len(segments) == 1
        assert segments[0].content == "Hello world"
        assert segments[0].metadata == {}

    def test_dispatch_unsupported_raises(self):
        with pytest.raises(ValueError, match="Unsupported"):
            parse_file_segments(b"data", "file.xyz")

    def test_empty_content_returns_empty(self):
        segments = parse_file_segments(b"", "test.txt")
        assert segments == []


class TestParsePdfSegments:
    def test_pdf_segments_per_page(self):
        """Each page of a PDF becomes a separate segment with page metadata."""
        import pymupdf

        doc = pymupdf.open()
        page1 = doc.new_page()
        page1.insert_text((72, 72), "Page one content")
        page2 = doc.new_page()
        page2.insert_text((72, 72), "Page two content")
        pdf_bytes = doc.tobytes()
        doc.close()

        segments = parse_file_segments(pdf_bytes, "report.pdf")
        assert len(segments) == 2
        assert "Page one content" in segments[0].content
        assert segments[0].metadata == {"page": 1}
        assert "Page two content" in segments[1].content
        assert segments[1].metadata == {"page": 2}


class TestParsePptxSegments:
    def test_pptx_segments_per_slide(self):
        from pptx import Presentation
        import io

        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])
        txBox = slide1.shapes.add_textbox(0, 0, 100, 100)
        txBox.text_frame.paragraphs[0].text = "Slide one"

        slide2 = prs.slides.add_slide(prs.slide_layouts[5])
        txBox2 = slide2.shapes.add_textbox(0, 0, 100, 100)
        txBox2.text_frame.paragraphs[0].text = "Slide two"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        segments = parse_file_segments(buf.read(), "deck.pptx")
        assert len(segments) == 2
        assert "Slide one" in segments[0].content
        assert segments[0].metadata == {"slide": 1}
        assert "Slide two" in segments[1].content
        assert segments[1].metadata == {"slide": 2}


class TestParseExcelSegments:
    def test_excel_segments_per_sheet(self):
        import openpyxl
        import io

        wb = openpyxl.Workbook()
        ws1 = wb.active
        ws1.title = "Sales"
        ws1.append(["Product", "Revenue"])
        ws1.append(["Widget", 100])

        ws2 = wb.create_sheet("Costs")
        ws2.append(["Item", "Amount"])
        ws2.append(["Rent", 500])

        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)

        segments = parse_file_segments(buf.read(), "data.xlsx")
        assert len(segments) == 2
        assert segments[0].metadata == {"sheet": "Sales"}
        assert segments[1].metadata == {"sheet": "Costs"}
        assert "Widget" in segments[0].content
        assert "Rent" in segments[1].content

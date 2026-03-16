"""Test data generators for integration tests."""

import io


def create_test_pdf(pages: list[str]) -> bytes:
    """Create a PDF with the given text on each page."""
    import pymupdf

    doc = pymupdf.open()
    for page_text in pages:
        page = doc.new_page()
        page.insert_text((72, 72), page_text, fontsize=12)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    return buf.getvalue()


def create_test_pptx(slides: list[str]) -> bytes:
    """Create a PPTX with one text box per slide."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for slide_text in slides:
        slide_layout = prs.slide_layouts[5]  # blank
        slide = prs.slides.add_slide(slide_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
        txBox.text_frame.text = slide_text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def create_test_excel(sheets: dict[str, list[list[str]]]) -> bytes:
    """Create an Excel file with named sheets.

    sheets: {"SheetName": [["header1", "header2"], ["val1", "val2"], ...]}
    """
    import openpyxl

    wb = openpyxl.Workbook()
    # Remove default sheet
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def create_test_docx(paragraphs: list[str]) -> bytes:
    """Create a Word (.docx) file with the given paragraphs."""
    import docx

    doc = docx.Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def create_test_csv(headers: list[str], rows: list[list[str]]) -> bytes:
    """Create a CSV file with the given headers and rows."""
    import csv

    buf = io.StringIO()
    writer = csv.writer(buf)
    writer.writerow(headers)
    writer.writerows(rows)
    return buf.getvalue().encode("utf-8")


def create_test_image(text: str, width: int = 400, height: int = 300) -> bytes:
    """Create a PNG image with the given text drawn on it."""
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Draw the text and some shapes to give the VL model something to describe
    draw.text((20, 20), text, fill=(0, 0, 0))

    # Draw some boxes to simulate an architecture diagram
    draw.rectangle([50, 80, 180, 140], outline=(0, 0, 200), width=2)
    draw.text((70, 100), "Web Server", fill=(0, 0, 200))

    draw.rectangle([220, 80, 350, 140], outline=(200, 0, 0), width=2)
    draw.text((240, 100), "Database", fill=(200, 0, 0))

    # Arrow between boxes
    draw.line([(180, 110), (220, 110)], fill=(0, 0, 0), width=2)
    draw.polygon([(215, 105), (220, 110), (215, 115)], fill=(0, 0, 0))

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()

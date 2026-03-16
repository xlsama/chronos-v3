"""File parsers for extracting text from various document formats."""

import csv
import io
import json
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class ParsedSegment:
    content: str
    metadata: dict = field(default_factory=dict)


IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp"}


def is_image(filename: str) -> bool:
    return Path(filename).suffix.lower() in IMAGE_EXTENSIONS


def parse_text(file_bytes: bytes, _filename: str = "") -> str:
    """Parse plain text / markdown files."""
    return file_bytes.decode("utf-8", errors="replace")


def parse_csv(file_bytes: bytes, _filename: str = "") -> str:
    """Parse CSV files into a readable text table."""
    text = file_bytes.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    if not rows:
        return ""
    # Format as markdown table
    header = rows[0]
    lines = ["| " + " | ".join(header) + " |"]
    lines.append("| " + " | ".join("---" for _ in header) + " |")
    for row in rows[1:]:
        # Pad row to match header length
        padded = row + [""] * (len(header) - len(row))
        lines.append("| " + " | ".join(padded[:len(header)]) + " |")
    return "\n".join(lines)


def parse_pdf(file_bytes: bytes, _filename: str = "") -> str:
    """Parse PDF files using pymupdf."""
    segments = parse_pdf_segments(file_bytes)
    return "\n\n".join(seg.content for seg in segments)


def parse_pdf_segments(file_bytes: bytes) -> list[ParsedSegment]:
    """Parse PDF, returning one segment per page."""
    import pymupdf

    doc = pymupdf.open(stream=file_bytes, filetype="pdf")
    segments = []
    for i, page in enumerate(doc, 1):
        text = page.get_text().strip()
        if text:
            segments.append(ParsedSegment(content=text, metadata={"page": i}))
    doc.close()
    return segments


def parse_docx(file_bytes: bytes, _filename: str = "") -> str:
    """Parse Word (.docx) files using python-docx."""
    import docx

    doc = docx.Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def parse_excel(file_bytes: bytes, filename: str = "") -> str:
    """Parse Excel (.xlsx/.xls) files using openpyxl."""
    segments = parse_excel_segments(file_bytes)
    return "\n\n".join(seg.content for seg in segments)


def parse_excel_segments(file_bytes: bytes) -> list[ParsedSegment]:
    """Parse Excel, returning one segment per sheet."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    segments = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        lines = [f"### Sheet: {sheet_name}"]
        header = [str(c) if c is not None else "" for c in rows[0]]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join("---" for _ in header) + " |")
        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            padded = cells + [""] * (len(header) - len(cells))
            lines.append("| " + " | ".join(padded[:len(header)]) + " |")
        segments.append(ParsedSegment(content="\n".join(lines), metadata={"sheet": sheet_name}))
    wb.close()
    return segments


def parse_pptx(file_bytes: bytes, _filename: str = "") -> str:
    """Parse PowerPoint (.pptx) files using python-pptx."""
    segments = parse_pptx_segments(file_bytes)
    return "\n\n".join(seg.content for seg in segments)


def parse_pptx_segments(file_bytes: bytes) -> list[ParsedSegment]:
    """Parse PPTX, returning one segment per slide."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    segments = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            content = f"### Slide {i}\n\n" + "\n".join(texts)
            segments.append(ParsedSegment(content=content, metadata={"slide": i}))
    return segments


def parse_html(file_bytes: bytes, _filename: str = "") -> str:
    """Parse HTML files using BeautifulSoup."""
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(file_bytes, "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n", strip=True)


def parse_json(file_bytes: bytes, _filename: str = "") -> str:
    """Parse JSON files into formatted text."""
    text = file_bytes.decode("utf-8", errors="replace")
    data = json.loads(text)
    return json.dumps(data, indent=2, ensure_ascii=False)


# Mapping of file extensions to parser functions
PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".xlsx": parse_excel,
    ".xls": parse_excel,
    ".csv": parse_csv,
    ".html": parse_html,
    ".htm": parse_html,
    ".json": parse_json,
    ".yaml": parse_text,
    ".yml": parse_text,
    ".log": parse_text,
    ".md": parse_text,
    ".txt": parse_text,
}

# Supported file extensions for upload (text + image)
SUPPORTED_EXTENSIONS = set(PARSERS.keys()) | IMAGE_EXTENSIONS

# Segment parsers for formats that have natural sub-document boundaries
SEGMENT_PARSERS = {
    ".pdf": lambda b, _fn: parse_pdf_segments(b),
    ".pptx": lambda b, _fn: parse_pptx_segments(b),
    ".xlsx": lambda b, _fn: parse_excel_segments(b),
    ".xls": lambda b, _fn: parse_excel_segments(b),
}


def parse_file_segments(file_bytes: bytes, filename: str) -> list[ParsedSegment]:
    """Parse a file into segments with metadata.

    For PDF/PPTX/Excel, returns per-page/slide/sheet segments.
    For other text formats, returns a single segment with empty metadata.

    Raises:
        ValueError: If the file format is not supported.
    """
    ext = Path(filename).suffix.lower()
    seg_parser = SEGMENT_PARSERS.get(ext)
    if seg_parser:
        return seg_parser(file_bytes, filename)

    parser = PARSERS.get(ext)
    if not parser:
        raise ValueError(
            f"Unsupported file format: {ext}. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    content = parser(file_bytes, filename)
    if not content.strip():
        return []
    return [ParsedSegment(content=content)]


def parse_file(file_bytes: bytes, filename: str) -> str:
    """Parse a file and extract text content.

    Args:
        file_bytes: Raw file bytes.
        filename: Original filename (used to determine format).

    Returns:
        Extracted text content.

    Raises:
        ValueError: If the file format is not supported.
    """
    segments = parse_file_segments(file_bytes, filename)
    return "\n\n".join(seg.content for seg in segments)
